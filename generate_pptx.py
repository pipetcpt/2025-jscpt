#!/usr/bin/env python3
"""
Generate PowerPoint presentation from NONMEM IRT Model Infographics
Enhanced version with detailed discussions and interpretations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(44, 62, 80)
DARK_SECONDARY = RGBColor(52, 73, 94)
BLUE_ACCENT = RGBColor(52, 152, 219)
GREEN_ACCENT = RGBColor(46, 204, 113)
RED_ACCENT = RGBColor(231, 76, 60)
ORANGE_ACCENT = RGBColor(243, 156, 18)
PURPLE_ACCENT = RGBColor(155, 89, 182)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(236, 240, 241)
MID_GRAY = RGBColor(189, 195, 199)
DIM_GRAY = RGBColor(149, 165, 166)


def add_title_slide(prs):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()

    # Decorative accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.8), Inches(0.15), Inches(2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = BLUE_ACCENT
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IRT Model for EDSS in Multiple Sclerosis"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Longitudinal Pharmacometric Analysis of Cladribine Treatment"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Reference
    ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.5))
    tf = ref_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Based on Novakovic et al., 2016"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Korean/Japanese text
    lang_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(1.5))
    tf = lang_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "다발성 경화증 EDSS 항목반응이론 모델 | 클라드리빈 치료 효과 분석"
    p.font.size = Pt(16)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "多発性硬化症 EDSS 項目反応理論モデル | クラドリビン治療効果分析"
    p2.font.size = Pt(16)
    p2.font.color.rgb = MID_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # Footer with key info
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.7))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "NONMEM 7.3 | Laplacian Conditional Estimation | Item Response Theory"
    p.font.size = Pt(14)
    p.font.color.rgb = DIM_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_agenda_slide(prs):
    """Add agenda/overview slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Presentation Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Agenda items
    agenda_items = [
        ("1", "Model Structure", "Longitudinal IRT framework with drug effects", BLUE_ACCENT),
        ("2", "Parameter Estimates", "Population PK/PD parameters and variability", GREEN_ACCENT),
        ("3", "Random Effects", "Correlation structure via Cholesky decomposition", PURPLE_ACCENT),
        ("4", "IRT Analysis", "Item characteristic curves and information functions", ORANGE_ACCENT),
        ("5", "Drug Effects", "Symptomatic and disease-modifying mechanisms", RED_ACCENT),
        ("6", "Clinical Implications", "Key insights for MS treatment optimization", BLUE_ACCENT),
    ]

    y_start = Inches(1.6)
    item_height = Inches(0.85)

    for i, (num, title, desc, color) in enumerate(agenda_items):
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y_start + i * item_height,
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.8), y_start + i * item_height + Inches(0.08), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_txt = slide.shapes.add_textbox(Inches(1.5), y_start + i * item_height, Inches(4), Inches(0.4))
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BG

        # Description
        desc_txt = slide.shapes.add_textbox(Inches(1.5), y_start + i * item_height + Inches(0.35), Inches(10), Inches(0.4))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_SECONDARY


def add_background_slide(prs):
    """Add clinical background slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinical Background: Multiple Sclerosis & EDSS"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column - MS Background
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_box.line.color.rgb = BLUE_ACCENT
    left_box.line.width = Pt(2)

    ms_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4))
    tf = ms_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Multiple Sclerosis (MS)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    ms_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(5.8), Inches(2.1))
    tf = ms_text.text_frame
    tf.word_wrap = True

    points = [
        "• Chronic autoimmune disease of the CNS",
        "• Affects ~2.8 million people worldwide",
        "• Progressive neurological disability over time",
        "• Characterized by relapses and remissions",
        "• Cladribine: oral disease-modifying therapy"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(6)

    # Right column - EDSS
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    right_box.line.color.rgb = GREEN_ACCENT
    right_box.line.width = Pt(2)

    edss_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4))
    tf = edss_title.text_frame
    p = tf.paragraphs[0]
    p.text = "EDSS (Expanded Disability Status Scale)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    edss_text = slide.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.8), Inches(2.1))
    tf = edss_text.text_frame
    tf.word_wrap = True

    points = [
        "• Gold standard for MS disability assessment",
        "• Composite score from 0 (normal) to 10 (death)",
        "• Based on 8 functional system subscores:",
        "  - Pyramidal, Cerebellar, Brainstem, Sensory",
        "  - Bowel/Bladder, Visual, Mental, Ambulation"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(6)

    # Bottom - Why IRT?
    bottom_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.8)
    )
    bottom_box.fill.solid()
    bottom_box.fill.fore_color.rgb = RGBColor(255, 248, 240)
    bottom_box.line.color.rgb = ORANGE_ACCENT
    bottom_box.line.width = Pt(2)

    irt_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4))
    tf = irt_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Why Item Response Theory (IRT) for EDSS?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE_ACCENT

    irt_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12.3), Inches(2))
    tf = irt_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Traditional Analysis Limitations:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BG

    limitations = [
        "• EDSS is ordinal, not continuous → standard regression inappropriate",
        "• Subscores have different scales (0-4, 0-5, 0-6, 0-9) → not directly comparable",
        "• Floor/ceiling effects mask true disability changes"
    ]

    for point in limitations:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SECONDARY
        p.space_before = Pt(3)

    p = tf.add_paragraph()
    p.text = "IRT Solution: Models a latent 'disability' trait that underlies all observed subscores, providing unified disease metric"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.space_before = Pt(10)


def add_image_slide_with_notes(prs, image_path, title, subtitle, notes, note_color=BLUE_ACCENT):
    """Add slide with image and discussion notes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.85))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(9), Inches(0.45))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.52), Inches(12.733), Inches(0.3))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = MID_GRAY

    # Add image (left side, smaller to make room for notes)
    if os.path.exists(image_path):
        img_width = Inches(8.5)
        img_left = Inches(0.2)
        img_top = Inches(0.95)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    # Notes panel on right
    notes_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.9), Inches(0.95), Inches(4.2), Inches(6.35)
    )
    notes_box.fill.solid()
    notes_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    notes_box.line.color.rgb = note_color
    notes_box.line.width = Pt(2)

    # Notes header
    notes_header = slide.shapes.add_textbox(Inches(9.1), Inches(1.05), Inches(3.8), Inches(0.4))
    tf = notes_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Points & Discussion"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = note_color

    # Notes content
    notes_content = slide.shapes.add_textbox(Inches(9.1), Inches(1.5), Inches(3.8), Inches(5.6))
    tf = notes_content.text_frame
    tf.word_wrap = True

    for i, note in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + note
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(8)
        p.space_after = Pt(4)


def add_model_structure_discussion(prs):
    """Add detailed model structure discussion slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Model Structure: Mathematical Framework"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Disease progression equation box
    eq_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.1), Inches(12.7), Inches(1.4)
    )
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = RGBColor(232, 244, 253)
    eq_box.line.color.rgb = BLUE_ACCENT
    eq_box.line.width = Pt(2)

    eq_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.35))
    tf = eq_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Disease Latent Variable Equation:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    eq_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8))
    tf = eq_text.text_frame
    p = tf.paragraphs[0]
    p.text = "PD = P₁ + (θ₁ + P₂) × (t/365)^θ₂ × (1 - Ef_prot) - Ef_symp"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Where: P₁ = baseline disability (η₁), P₂ = individual slope (η₂), θ₁ = population slope, θ₂ = power parameter"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Three component boxes
    components = [
        ("Disease Progression", "Power Model",
         ["Slope (θ₁) = 0.093", "Power (θ₂) = 0.710", "Non-linear time course", "Sub-linear progression pattern"],
         BLUE_ACCENT, Inches(0.3)),
        ("Symptomatic Effect", "Emax Model",
         ["Emax = 0.17", "EC50 = 408.29", "Exposure-dependent", "Immediate effect on symptoms"],
         GREEN_ACCENT, Inches(4.5)),
        ("Protective Effect", "Disease-Modifying",
         ["Effect = 20.9%", "Exposure-independent", "Slows progression rate", "Long-term benefit"],
         RED_ACCENT, Inches(8.7)),
    ]

    for title, model_type, points, color, x_pos in components:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(2.7), Inches(4.0), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Header
        header = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(2.8), Inches(3.7), Inches(0.35))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        subhead = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(3.1), Inches(3.7), Inches(0.25))
        tf = subhead.text_frame
        p = tf.paragraphs[0]
        p.text = model_type
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = DARK_SECONDARY

        content = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(3.4), Inches(3.7), Inches(1.6))
        tf = content.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_BG
            p.space_before = Pt(4)

    # IRT link explanation
    irt_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(5.4), Inches(12.7), Inches(1.8)
    )
    irt_box.fill.solid()
    irt_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    irt_box.line.color.rgb = ORANGE_ACCENT
    irt_box.line.width = Pt(2)

    irt_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.35))
    tf = irt_title.text_frame
    p = tf.paragraphs[0]
    p.text = "IRT Link: Connecting Latent Disease to Observed Subscores"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE_ACCENT

    irt_content = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2))
    tf = irt_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Graded Response Model:  P(X ≥ k | θ) = 1 / (1 + exp(-a × (θ - bₖ)))"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BG

    p2 = tf.add_paragraph()
    p2.text = "• Each functional system has unique discrimination (a) and boundary (b) parameters"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SECONDARY
    p2.space_before = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "• Higher discrimination → item better differentiates disability levels | Boundaries define score thresholds"
    p3.font.size = Pt(12)
    p3.font.color.rgb = DARK_SECONDARY


def add_correlation_discussion(prs):
    """Add correlation interpretation slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Random Effects: Clinical Interpretation of Correlations"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Cholesky explanation
    chol_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.1), Inches(6.2), Inches(1.5)
    )
    chol_box.fill.solid()
    chol_box.fill.fore_color.rgb = RGBColor(245, 245, 255)
    chol_box.line.color.rgb = PURPLE_ACCENT
    chol_box.line.width = Pt(2)

    chol_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.35))
    tf = chol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Cholesky Parameterization"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PURPLE_ACCENT

    chol_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.9))
    tf = chol_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Ensures positive-definite covariance matrix"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_BG
    p2 = tf.add_paragraph()
    p2.text = "• Correlations estimated as separate parameters"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_BG
    p3 = tf.add_paragraph()
    p3.text = "• FREM: Covariates as random effects for correlation"
    p3.font.size = Pt(12)
    p3.font.color.rgb = DARK_BG

    # Key correlations
    corr_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.1), Inches(6.2), Inches(1.5)
    )
    corr_box.fill.solid()
    corr_box.fill.fore_color.rgb = RGBColor(240, 255, 245)
    corr_box.line.color.rgb = GREEN_ACCENT
    corr_box.line.width = Pt(2)

    corr_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35))
    tf = corr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "5×5 Correlation Structure"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    corr_text = slide.shapes.add_textbox(Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.9))
    tf = corr_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Random effects: η₁ (Disability), η₂ (Slope),"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_BG
    p2 = tf.add_paragraph()
    p2.text = "η₃ (Age), η₄ (MSD), η₅ (EXNB)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_BG
    p3 = tf.add_paragraph()
    p3.text = "Plus η₆ (Drug Emax) - separate block"
    p3.font.size = Pt(12)
    p3.font.color.rgb = DARK_BG

    # Clinical interpretations
    interp_data = [
        ("Age ↔ MS Duration", "r = 0.458", "Moderate Positive",
         "Older patients have longer disease duration on average. Expected clinical finding validates model.",
         BLUE_ACCENT),
        ("Disability ↔ MS Duration", "r = 0.273", "Weak Positive",
         "Longer disease associates with higher baseline disability. Supports progressive nature of MS.",
         GREEN_ACCENT),
        ("Disability ↔ Age", "r = 0.265", "Weak Positive",
         "Age contributes to disability beyond disease duration. May reflect age-related neural reserve decline.",
         ORANGE_ACCENT),
        ("MS Duration ↔ EXNB", "r = -0.115", "Weak Negative",
         "Longer disease duration associated with fewer recent exacerbations. Possible transition to progressive phase.",
         RED_ACCENT),
    ]

    y_start = Inches(2.8)
    box_height = Inches(1.05)

    for i, (pair, corr, strength, interp, color) in enumerate(interp_data):
        row = i // 2
        col = i % 2
        x_pos = Inches(0.3) + col * Inches(6.5)
        y_pos = y_start + row * (box_height + Inches(0.15))

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(6.2), box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Header line
        header = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.08), Inches(3.5), Inches(0.3))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = pair
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        corr_val = slide.shapes.add_textbox(x_pos + Inches(3.8), y_pos + Inches(0.08), Inches(2.2), Inches(0.3))
        tf = corr_val.text_frame
        p = tf.paragraphs[0]
        p.text = f"{corr} ({strength})"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_SECONDARY
        p.alignment = PP_ALIGN.RIGHT

        # Interpretation
        interp_txt = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.4), Inches(5.9), Inches(0.6))
        tf = interp_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = interp
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BG

    # Bottom note
    note_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Note: Weak correlations (|r| < 0.3) suggest relatively independent random effects, supporting model identifiability"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = DIM_GRAY


def add_irt_interpretation_slide(prs):
    """Add IRT clinical interpretation slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IRT Analysis: Clinical Implications of Item Parameters"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # High discrimination items
    high_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.1), Inches(6.2), Inches(2.6)
    )
    high_box.fill.solid()
    high_box.fill.fore_color.rgb = RGBColor(232, 255, 232)
    high_box.line.color.rgb = GREEN_ACCENT
    high_box.line.width = Pt(2)

    high_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.35))
    tf = high_title.text_frame
    p = tf.paragraphs[0]
    p.text = "High Discrimination Items (slope > 2.5)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    high_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(2))
    tf = high_content.text_frame
    tf.word_wrap = True

    high_items = [
        ("Ambulation (3.64)", "Most sensitive to disability changes. Walking ability strongly reflects underlying disease severity."),
        ("Pyramidal (3.17)", "Motor function highly informative. Corticospinal tract involvement is a key MS feature."),
        ("Cerebellar (2.87)", "Balance and coordination discriminate well between disability levels."),
    ]

    for i, (item, desc) in enumerate(high_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_SECONDARY

    # Low discrimination items
    low_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.1), Inches(6.2), Inches(2.6)
    )
    low_box.fill.solid()
    low_box.fill.fore_color.rgb = RGBColor(255, 245, 238)
    low_box.line.color.rgb = ORANGE_ACCENT
    low_box.line.width = Pt(2)

    low_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35))
    tf = low_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Lower Discrimination Items (slope < 1.5)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE_ACCENT

    low_content = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.8), Inches(2))
    tf = low_content.text_frame
    tf.word_wrap = True

    low_items = [
        ("Visual (0.44)", "Optic neuritis often occurs early and may recover. Less reflective of overall progression."),
        ("Mental (0.91)", "Cognitive symptoms variable and may not correlate linearly with physical disability."),
        ("Sensory (0.99)", "Sensory symptoms fluctuate and are subjective. Less reliable disability indicator."),
    ]

    for i, (item, desc) in enumerate(low_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_SECONDARY

    # Clinical implications box
    impl_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(3.9), Inches(12.7), Inches(3.2)
    )
    impl_box.fill.solid()
    impl_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    impl_box.line.color.rgb = BLUE_ACCENT
    impl_box.line.width = Pt(2)

    impl_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.35))
    tf = impl_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinical & Trial Design Implications"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    implications = [
        ("Trial Endpoints", "Consider weighting functional systems by discrimination. Ambulation-focused endpoints may be more sensitive to treatment effects."),
        ("Patient Monitoring", "Track high-discrimination items more closely for early detection of progression or treatment response."),
        ("Sample Size", "IRT-based analysis can increase power by modeling the latent trait directly, potentially reducing required sample size."),
        ("Precision Medicine", "Patients with primarily visual or sensory symptoms may show less response on standard EDSS despite treatment benefit."),
    ]

    impl_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6))
    tf = impl_content.text_frame
    tf.word_wrap = True

    for i, (topic, desc) in enumerate(implications):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {topic}: "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(8)

        # Add description as run in same paragraph
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = DARK_SECONDARY


def add_drug_effect_discussion(prs):
    """Add drug effect mechanism discussion slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cladribine: Dual Mechanism Drug Effects"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Symptomatic effect box
    symp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.1), Inches(6.2), Inches(3.0)
    )
    symp_box.fill.solid()
    symp_box.fill.fore_color.rgb = RGBColor(232, 255, 243)
    symp_box.line.color.rgb = GREEN_ACCENT
    symp_box.line.width = Pt(3)

    symp_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4))
    tf = symp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Symptomatic Effect"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    symp_model = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.4))
    tf = symp_model.text_frame
    p = tf.paragraphs[0]
    p.text = "Ef_symp = Emax × Exposure / (Exposure + EC50)"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BG
    p.font.bold = True

    symp_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(1.9))
    tf = symp_content.text_frame
    tf.word_wrap = True

    symp_points = [
        "Emax = 0.17 (17% maximum reduction)",
        "EC50 = 408.29 (half-maximal exposure)",
        "Exposure-dependent: higher dose → greater effect",
        "Direct reduction in disability latent variable",
        "Rapid onset, dependent on drug levels",
        "Surrogate exposure: CD × 104.5 / CrCL"
    ]

    for i, point in enumerate(symp_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(5)

    # Protective effect box
    prot_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.1), Inches(6.2), Inches(3.0)
    )
    prot_box.fill.solid()
    prot_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    prot_box.line.color.rgb = RED_ACCENT
    prot_box.line.width = Pt(3)

    prot_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4))
    tf = prot_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Protective (Disease-Modifying) Effect"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT

    prot_model = slide.shapes.add_textbox(Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.4))
    tf = prot_model.text_frame
    p = tf.paragraphs[0]
    p.text = "Progression × (1 - Ef_prot) where Ef_prot = 0.209"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BG
    p.font.bold = True

    prot_content = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.8), Inches(1.9))
    tf = prot_content.text_frame
    tf.word_wrap = True

    prot_points = [
        "20.9% reduction in progression rate",
        "Exposure-INDEPENDENT (binary effect)",
        "Slows the underlying disease process",
        "Cumulative benefit over time",
        "Active when TRT ≥ 1 and time > 0",
        "Mechanism: lymphocyte depletion"
    ]

    for i, point in enumerate(prot_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(5)

    # Clinical significance box
    clin_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.85)
    )
    clin_box.fill.solid()
    clin_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    clin_box.line.color.rgb = BLUE_ACCENT
    clin_box.line.width = Pt(2)

    clin_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4))
    tf = clin_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinical Significance & Therapeutic Implications"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    clin_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.2))
    tf = clin_content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Dual Mechanism Advantage:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BG

    clin_points = [
        "Symptomatic effect provides immediate benefit patients can perceive (improved daily function)",
        "Protective effect provides long-term disease modification (delayed disability progression)",
        "Combined effect: both symptom relief AND slower progression trajectory",
        "Model enables prediction of long-term outcomes from short-term exposure data",
    ]

    for point in clin_points:
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SECONDARY
        p.space_before = Pt(5)

    p = tf.add_paragraph()
    p.text = "Key Insight: The protective effect (20.9%) compounds over time - a patient treated for 5 years accumulates substantial disability reduction vs. placebo"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_before = Pt(10)


def add_summary_slide(prs):
    """Add comprehensive summary slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary: Key Findings & Clinical Implications"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Key findings grid
    findings = [
        ("Disease Model", "Power function progression\nSlope = 0.093, Power = 0.71\nNon-linear, sub-linear course", BLUE_ACCENT),
        ("Drug Effects", "Dual mechanism: Symptomatic (Emax)\n+ Protective (20.9% rate reduction)\nBoth contribute to efficacy", GREEN_ACCENT),
        ("IRT Insights", "Ambulation most informative (a=3.64)\nVisual least informative (a=0.44)\nConsider weighted endpoints", ORANGE_ACCENT),
        ("Correlations", "Age-MSD: r=0.46 (expected)\nWeak inter-correlations support\nmodel identifiability", PURPLE_ACCENT),
        ("Population", "Mean age: 38.6 years\nMean MS duration: 8.7 years\nMean exacerbations: 1.35", RED_ACCENT),
        ("Model Fit", "OFV = 418.626\nLaplacian estimation\n241 obs, 3 subjects (simulated)", BLUE_ACCENT),
    ]

    for i, (title, content, color) in enumerate(findings):
        row = i // 3
        col = i % 3
        x_pos = Inches(0.4) + col * Inches(4.3)
        y_pos = Inches(1.2) + row * Inches(2.5)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(4.0), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_SECONDARY
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # Title
        title_txt = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.1), Inches(3.7), Inches(0.4))
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # Content
        content_txt = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.55), Inches(3.7), Inches(1.5))
        tf = content_txt.text_frame
        tf.word_wrap = True
        for j, line in enumerate(content.split('\n')):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_GRAY
            p.space_before = Pt(4)

    # Bottom message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.8))
    tf = msg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Conclusion: IRT-based longitudinal modeling provides mechanistic insights into cladribine's dual action,"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "enabling optimized dosing strategies and more sensitive clinical trial endpoints for MS treatment."
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER


def add_references_slide(prs):
    """Add references slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "References & Further Reading"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # References
    refs = [
        "Primary Reference:",
        "• Novakovic AM, et al. (2016). Longitudinal Item Response Theory Model for EDSS in Multiple Sclerosis. CPT: Pharmacometrics & Systems Pharmacology.",
        "",
        "IRT Methodology:",
        "• Samejima F. (1969). Estimation of latent ability using a response pattern of graded scores. Psychometrika Monograph Supplement, 34(4).",
        "• Ueckert S, et al. (2014). Modeling composite assessment data using item response theory. CPT: Pharmacometrics & Systems Pharmacology.",
        "",
        "NONMEM & Pharmacometrics:",
        "• Beal S, Sheiner LB, Boeckmann A, & Bauer RJ. (2009). NONMEM User's Guides. Icon Development Solutions.",
        "• Karlsson MO & Holford N. (2008). A Tutorial on Visual Predictive Checks. PAGE Meeting.",
        "",
        "Multiple Sclerosis & EDSS:",
        "• Kurtzke JF. (1983). Rating neurologic impairment in multiple sclerosis: an expanded disability status scale (EDSS). Neurology.",
        "• Thompson AJ, et al. (2018). Diagnosis of multiple sclerosis: 2017 revisions of the McDonald criteria. Lancet Neurology.",
        "",
        "Cladribine:",
        "• Giovannoni G, et al. (2010). A placebo-controlled trial of oral cladribine for relapsing multiple sclerosis. NEJM.",
    ]

    ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(6.2))
    tf = ref_box.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(refs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref

        if ref.endswith(':') and not ref.startswith('•'):
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = BLUE_ACCENT
            p.space_before = Pt(12)
        elif ref == "":
            p.font.size = Pt(8)
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_BG
            p.space_before = Pt(3)


def main():
    """Generate PowerPoint presentation"""
    print("=" * 60)
    print("Generating Enhanced PowerPoint Presentation")
    print("with Detailed Discussions and Interpretations")
    print("=" * 60)
    print()

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    print("Creating slide 1: Title...")
    add_title_slide(prs)

    # Slide 2: Agenda
    print("Creating slide 2: Agenda...")
    add_agenda_slide(prs)

    # Slide 3: Clinical Background
    print("Creating slide 3: Clinical Background...")
    add_background_slide(prs)

    # Slide 4: Model Structure Image + Notes
    print("Creating slide 4: Model Structure (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_01_model_structure.png",
        "Model Structure Overview",
        "Longitudinal IRT Model Architecture",
        [
            "Central latent variable (PD) represents true underlying disability",
            "Power model captures non-linear disease progression over time",
            "Two distinct drug effect mechanisms: symptomatic and protective",
            "FREM approach incorporates Age, MS Duration, Exacerbation count as covariates",
            "8 functional systems linked via IRT graded response model",
            "Cholesky decomposition ensures valid correlation structure",
        ],
        BLUE_ACCENT
    )

    # Slide 5: Model Structure Discussion
    print("Creating slide 5: Model Structure Discussion...")
    add_model_structure_discussion(prs)

    # Slide 6: Parameter Estimates Image + Notes
    print("Creating slide 6: Parameter Estimates (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_02_parameters.png",
        "Parameter Estimates Summary",
        "Population Parameters and Variability",
        [
            "Disease progression slope (0.093) indicates slow accumulation of disability",
            "Power < 1 (0.71) suggests decelerating progression over time",
            "Large variance in Age (99.3) reflects patient heterogeneity",
            "Ambulation has highest discrimination (3.64) - most informative subscore",
            "Visual has lowest discrimination (0.44) - less reliable indicator",
            "OFV = 418.626 for model evaluation",
        ],
        GREEN_ACCENT
    )

    # Slide 7: Correlation Matrix Image + Notes
    print("Creating slide 7: Correlation Matrix (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_03_correlation_matrix.png",
        "Random Effects Correlation Structure",
        "Cholesky Parameterization of 5×5 Omega Matrix",
        [
            "Age-MSD correlation (0.46) is clinically expected and validates the model",
            "Disability correlates weakly with Age (0.26) and MSD (0.27)",
            "Slope shows minimal correlation with covariates (good for identifiability)",
            "Negative EXNB correlations may reflect disease phase transitions",
            "Overall weak correlations support parameter estimability",
            "FREM enables covariate effects without traditional covariate modeling",
        ],
        PURPLE_ACCENT
    )

    # Slide 8: Correlation Discussion
    print("Creating slide 8: Correlation Discussion...")
    add_correlation_discussion(prs)

    # Slide 9: ICC Curves Image + Notes
    print("Creating slide 9: ICC Curves (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_04_icc_curves.png",
        "Item Characteristic Curves",
        "Probability of Scoring ≥ k at Each Disability Level",
        [
            "Steeper curves = higher discrimination = better differentiation",
            "Ambulation shows sharp transitions - excellent disability marker",
            "Visual shows gradual curves - less sensitive to changes",
            "Curve spacing indicates boundary (threshold) locations",
            "Items most informative where curves are steepest",
            "Can guide selection of endpoints for clinical trials",
        ],
        ORANGE_ACCENT
    )

    # Slide 10: IRT Interpretation
    print("Creating slide 10: IRT Interpretation...")
    add_irt_interpretation_slide(prs)

    # Slide 11: Drug Effects Image + Notes
    print("Creating slide 11: Drug Effects (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_05_drug_effects.png",
        "Cladribine Drug Effects",
        "Symptomatic and Disease-Modifying Mechanisms",
        [
            "Emax model: saturable symptomatic effect with exposure",
            "EC50 = 408.29 guides dosing for optimal symptomatic relief",
            "Protective effect (20.9%) acts on progression rate, not absolute level",
            "Combined effects separate over time - greatest benefit long-term",
            "Model supports both immediate benefit and disease modification claims",
            "Enables simulation of various dosing regimens",
        ],
        RED_ACCENT
    )

    # Slide 12: Drug Effect Discussion
    print("Creating slide 12: Drug Effect Discussion...")
    add_drug_effect_discussion(prs)

    # Slide 13: IRT Summary Image + Notes
    print("Creating slide 13: IRT Summary (with notes)...")
    add_image_slide_with_notes(
        prs,
        "infographic_06_irt_summary.png",
        "IRT Item Parameter Analysis",
        "Discrimination, Information, and Boundary Thresholds",
        [
            "Test Information Function peaks around θ = 2-3 (moderate disability)",
            "Less precision at extremes (very low or very high disability)",
            "Ambulation boundaries span wide θ range (0-9 score range)",
            "Item information curves show where each subscore is most useful",
            "Standard error inversely related to information",
            "Results support IRT as superior to simple sum scoring",
        ],
        BLUE_ACCENT
    )

    # Slide 14: Summary
    print("Creating slide 14: Summary...")
    add_summary_slide(prs)

    # Slide 15: References
    print("Creating slide 15: References...")
    add_references_slide(prs)

    # Save presentation
    output_file = "NONMEM_IRT_Model_Infographics.pptx"
    prs.save(output_file)

    print()
    print("=" * 60)
    print(f"PowerPoint presentation saved: {output_file}")
    print("=" * 60)
    print()
    print("Slides included (15 total):")
    print("  1.  Title Slide")
    print("  2.  Presentation Overview / Agenda")
    print("  3.  Clinical Background: MS & EDSS")
    print("  4.  Model Structure (Image + Key Points)")
    print("  5.  Model Structure: Mathematical Framework")
    print("  6.  Parameter Estimates (Image + Key Points)")
    print("  7.  Correlation Matrix (Image + Key Points)")
    print("  8.  Correlation: Clinical Interpretation")
    print("  9.  ICC Curves (Image + Key Points)")
    print("  10. IRT Analysis: Clinical Implications")
    print("  11. Drug Effects (Image + Key Points)")
    print("  12. Dual Mechanism Drug Effects Discussion")
    print("  13. IRT Summary (Image + Key Points)")
    print("  14. Summary: Key Findings")
    print("  15. References & Further Reading")


if __name__ == "__main__":
    main()
